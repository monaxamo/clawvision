#!/usr/bin/env python3
"""Render a ClawVision summary as HTML, PNG, Markdown, or PowerPoint."""

import argparse
import json
import re
import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt


TEMPLATE = r"""<!DOCTYPE html>
<html lang="{{lang}}">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{{title}}</title>
<style>
:root{--bg:#f5f6f8;--card:#fff;--text:#1a1a1a;--muted:#666;--accent:#2a9df4;--green:#4cd964;--orange:#ff9500;--red:#ff3b30;--border:#e5e7eb;}
:root.dark{--bg:#0f1115;--card:#1a1d23;--text:#e8e8e8;--muted:#9aa0a6;--border:#2c3038;--accent:#4aa8ff;--green:#5dd877;--orange:#ffae33;--red:#ff6659}
*{box-sizing:border-box}
body{margin:0;font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,Arial,sans-serif;background:var(--bg);color:var(--text);line-height:1.45;transition:background .2s,color .2s}
.wrap{max-width:760px;margin:0 auto;padding:20px}
.topbar{display:flex;justify-content:space-between;align-items:center;gap:12px;flex-wrap:wrap;margin-bottom:18px}
.controls{display:flex;gap:8px;flex-wrap:wrap}
.export-btn,.lang-btn,.theme-btn{border:none;border-radius:20px;padding:7px 13px;font-size:12px;font-weight:600;cursor:pointer;transition:.15s}
.export-btn{background:var(--accent);color:#fff}
.export-btn:hover{filter:brightness(1.1)}
.lang-switch,.theme-switch{display:flex;gap:6px;background:var(--card);border:1px solid var(--border);border-radius:20px;padding:4px}
.lang-btn,.theme-btn{background:transparent;color:var(--muted)}
.lang-btn.active,.theme-btn.active{background:var(--accent);color:#fff}
header{margin-bottom:18px}
header h1{font-size:22px;font-weight:700;margin:0 0 4px}
header p{color:var(--muted);margin:0;font-size:14px}
.badge{display:inline-block;background:#eef6ff;color:var(--accent);font-size:12px;font-weight:600;padding:4px 10px;border-radius:12px;margin-bottom:12px}
.tabs{display:flex;gap:8px;flex-wrap:wrap;margin-bottom:18px}
.tab{background:var(--card);border:1px solid var(--border);border-radius:20px;padding:8px 14px;font-size:13px;cursor:pointer;transition:.15s}
.tab.active{background:var(--accent);color:#fff;border-color:var(--accent)}
.panel{display:none;background:var(--card);border-radius:16px;padding:18px;box-shadow:0 2px 8px rgba(0,0,0,.04)}
.panel.active{display:block}
.lead{color:var(--muted);font-size:15px;margin-bottom:16px}
.flow{display:flex;align-items:center;gap:10px;flex-wrap:wrap;margin:16px 0}
.flow-item{background:#f8f9fa;border:1px solid var(--border);border-radius:12px;padding:10px 14px;min-width:110px;text-align:center}
.flow-item strong{display:block;font-size:14px;margin-bottom:2px}
.flow-item small{color:var(--muted);font-size:12px}
.arrow{color:var(--muted);font-size:18px}
.grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:12px;margin-top:12px}
.card{background:#f8f9fa;border-radius:14px;padding:14px}
.dark .flow-item,.dark .card{background:#22262e}
.card h3{margin:0 0 6px;font-size:16px}
.card p{margin:0;color:var(--muted);font-size:13px}
.checklist{margin:0;padding:0;list-style:none}
.checklist li{display:flex;justify-content:space-between;align-items:center;padding:10px 0;border-bottom:1px solid var(--border);font-size:14px}
.checklist li:last-child{border-bottom:none}
.status{font-size:12px;font-weight:600;color:var(--green)}
.status.pending{color:var(--orange)}
.status.blocked{color:var(--red)}
.two{display:grid;grid-template-columns:1fr 1fr;gap:12px}
ul.clean{padding-left:18px;margin:0}
ul.clean li{margin-bottom:6px;font-size:13px}
@media(max-width:600px){.two{grid-template-columns:1fr}.topbar{flex-direction:column;align-items:flex-start}}
</style>
</head>
<body>
<div class="wrap">
  <div class="topbar">
    <div class="controls">
      <a class="export-btn" href="{{md_file}}" download>Export Markdown</a>
      <a class="export-btn" href="{{pptx_file}}" download>Export PowerPoint</a>
    </div>
    <div class="controls">
      <div class="theme-switch"><button class="theme-btn" data-theme="light">☀</button><button class="theme-btn" data-theme="dark">🌙</button></div>
      <div class="lang-switch"><button class="lang-btn" data-lang="en">EN</button><button class="lang-btn" data-lang="ru">RU</button><button class="lang-btn" data-lang="zh">中文</button></div>
    </div>
  </div>

  <span class="badge">{{badge}}</span>
  <header>
    <h1>{{title}}</h1>
    <p>{{subtitle}}</p>
  </header>

  <div class="tabs">
    <div class="tab active" data-tab="main"><span id="main-title">{{tab_main}}</span></div>
    <div class="tab" data-tab="format"><span id="format-title">{{tab_format}}</span></div>
    <div class="tab" data-tab="built"><span id="built-title">{{tab_built}}</span></div>
    <div class="tab" data-tab="next"><span id="next-title">{{tab_next}}</span></div>
  </div>

  <div id="main" class="panel active">
    <p class="lead">{{main_takeaway}}</p>
    <div class="flow">{{flow_items}}</div>
    <div class="grid">{{metric_cards}}</div>
  </div>

  <div id="format" class="panel">
    <p class="lead">{{format_takeaway}}</p>
    <div class="two">
      <div class="card">
        <h3 id="dos-title">{{dos_title}}</h3>
        <ul class="clean">{{dos}}</ul>
      </div>
      <div class="card">
        <h3 id="donts-title">{{donts_title}}</h3>
        <ul class="clean">{{donts}}</ul>
      </div>
    </div>
  </div>

  <div id="built" class="panel">
    <ul class="checklist">{{checklist}}</ul>
  </div>

  <div id="next" class="panel">
    <p class="lead">{{next_takeaway}}</p>
    <ul class="clean">{{next_steps}}</ul>
  </div>
</div>
<script>
const LANGS={{langs_json}};
function setLang(l){
  document.body.dataset.lang=l;
  localStorage.setItem('cv-lang',l);
  document.querySelectorAll('.lang-btn').forEach(b=>b.classList.toggle('active',b.dataset.lang===l));
  const labels=LANGS[l]||LANGS.en;
  ['main','format','built','next'].forEach(k=>document.getElementById(k+'-title').textContent=labels[k]);
  document.getElementById('dos-title').textContent=labels.dos;
  document.getElementById('donts-title').textContent=labels.donts;
}
function setTheme(t){
  document.documentElement.classList.toggle('dark',t==='dark');
  localStorage.setItem('cv-theme',t);
  document.querySelectorAll('.theme-btn').forEach(b=>b.classList.toggle('active',b.dataset.theme===t));
}
document.querySelectorAll('.tab').forEach(t=>t.addEventListener('click',()=>{
  document.querySelectorAll('.tab').forEach(x=>x.classList.remove('active'));
  document.querySelectorAll('.panel').forEach(x=>x.classList.remove('active'));
  t.classList.add('active');
  document.getElementById(t.dataset.tab).classList.add('active');
}));
document.querySelectorAll('.lang-btn').forEach(b=>b.addEventListener('click',()=>setLang(b.dataset.lang)));
document.querySelectorAll('.theme-btn').forEach(b=>b.addEventListener('click',()=>setTheme(b.dataset.theme)));
setLang(localStorage.getItem('cv-lang')||'{{lang}}');
setTheme(localStorage.getItem('cv-theme')||'light');
</script>
</body>
</html>
"""


LANG_LABELS = {
    "en": {"main": "Main takeaway", "format": "Format", "built": "What we built", "next": "Next steps", "dos": "Do", "donts": "Don't", "export_md": "Export Markdown", "export_pptx": "Export PowerPoint", "badge": "ClawVision · Summary"},
    "ru": {"main": "Главный вывод", "format": "Формат", "built": "Что построено", "next": "Что дальше", "dos": "Нормально", "donts": "Риски", "export_md": "Экспорт Markdown", "export_pptx": "Экспорт PowerPoint", "badge": "ClawVision · Сводка"},
    "zh": {"main": "主要结论", "format": "形式", "built": "已完成", "next": "下一步", "dos": "建议", "donts": "风险", "export_md": "导出 Markdown", "export_pptx": "导出 PowerPoint", "badge": "ClawVision · 摘要"},
}


def _h(text: str) -> str:
    """Minimal HTML escape."""
    return (
        (text or "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def render_html(summary: dict, slug: str, md_file: str, pptx_file: str, lang: str = "en") -> str:
    labels = LANG_LABELS.get(lang, LANG_LABELS["en"])

    html = TEMPLATE
    html = html.replace("{{lang}}", _h(lang))
    html = html.replace("{{langs_json}}", json.dumps(LANG_LABELS, ensure_ascii=False))
    html = html.replace("{{md_file}}", _h(md_file))
    html = html.replace("{{pptx_file}}", _h(pptx_file))
    html = html.replace("{{title}}", _h(summary.get("title", "ClawVision summary")))
    html = html.replace("{{subtitle}}", _h(summary.get("subtitle", "")))
    html = html.replace("{{badge}}", _h(labels["badge"]))
    html = html.replace("{{tab_main}}", labels["main"])
    html = html.replace("{{tab_format}}", labels["format"])
    html = html.replace("{{tab_built}}", labels["built"])
    html = html.replace("{{tab_next}}", labels["next"])
    html = html.replace("{{main_takeaway}}", _h(summary.get("main_takeaway", "")))
    html = html.replace("{{format_takeaway}}", _h(summary.get("format_takeaway", "")))
    html = html.replace("{{next_takeaway}}", _h(summary.get("next_takeaway", "")))
    html = html.replace("{{dos_title}}", labels["dos"])
    html = html.replace("{{donts_title}}", labels["donts"])

    flow_items = ""
    for item in summary.get("flow", []):
        label = item.get("label", "")
        sub = item.get("sub", "")
        if label in ("→", "->"):
            flow_items += f'<div class="arrow">{_h(label)}</div>\n'
        else:
            flow_items += (
                f'<div class="flow-item"><strong>{_h(label)}</strong>'
                f'<small>{_h(sub)}</small></div>\n'
            )
    html = html.replace("{{flow_items}}", flow_items)

    metric_cards = ""
    for m in summary.get("metrics", []):
        metric_cards += (
            f'<div class="card"><h3>{_h(m.get("title", ""))}</h3>'
            f'<p>{_h(m.get("text", ""))}</p></div>\n'
        )
    html = html.replace("{{metric_cards}}", metric_cards)

    dos = "".join(f"<li>{_h(d)}</li>\n" for d in summary.get("dos", []))
    html = html.replace("{{dos}}", dos)

    donts = "".join(f"<li>{_h(d)}</li>\n" for d in summary.get("donts", []))
    html = html.replace("{{donts}}", donts)

    checklist = ""
    for c in summary.get("checklist", []):
        status = c.get("status", "pending")
        checklist += (
            f'<li>{_h(c.get("text", ""))} '
            f'<span class="status {status}">{_h(status)}</span></li>\n'
        )
    html = html.replace("{{checklist}}", checklist)

    next_steps = "".join(f"<li>{_h(s)}</li>\n" for s in summary.get("next_steps", []))
    html = html.replace("{{next_steps}}", next_steps)

    return html


def render_md(summary: dict, lang: str = "en") -> str:
    labels = LANG_LABELS.get(lang, LANG_LABELS["en"])
    lines = [
        f"# {summary.get('title', 'ClawVision summary')}",
        f"_{summary.get('subtitle', '')}_\n",
        f"## {labels['main']}",
        summary.get("main_takeaway", ""),
        "",
        f"## {labels['format']}",
        summary.get("format_takeaway", ""),
        "",
        f"### {labels['dos']}",
    ]
    lines.extend(f"- {d}" for d in summary.get("dos", []))
    lines.append("")
    lines.append(f"### {labels['donts']}")
    lines.extend(f"- {d}" for d in summary.get("donts", []))
    lines.append("")
    lines.append(f"## {labels['built']}")
    for c in summary.get("checklist", []):
        mark = "[x]" if c.get("status") == "ready" else "[ ]"
        lines.append(f"- {mark} {c.get('text', '')} ({c.get('status', 'pending')})")
    lines.append("")
    lines.append(f"## {labels['next']}")
    lines.append(summary.get("next_takeaway", ""))
    lines.append("")
    lines.extend(f"- {s}" for s in summary.get("next_steps", []))
    return "\n".join(lines)


def render_pptx(summary: dict, pptx_path: Path, lang: str = "en"):
    labels = LANG_LABELS.get(lang, LANG_LABELS["en"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(title, bullets):
        slide = prs.slides.add_slide(blank)
        t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        b = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(5.8))
        tf = b.text_frame
        tf.word_wrap = True
        for i, text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(18)
            p.space_after = Pt(10)
        return slide

    add_slide(summary.get("title", "ClawVision summary"), [summary.get("subtitle", "")])
    add_slide(labels["main"], [summary.get("main_takeaway", "")])
    add_slide(labels["format"], [summary.get("format_takeaway", "")])
    add_slide(labels["dos"], summary.get("dos", []))
    add_slide(labels["donts"], summary.get("donts", []))
    add_slide(labels["built"], [f"{c.get('text', '')} — {c.get('status', 'pending')}" for c in summary.get("checklist", [])])
    add_slide(labels["next"], [summary.get("next_takeaway", "")] + summary.get("next_steps", []))
    prs.save(str(pptx_path))


def screenshot_tabs(html_path: Path, slug: str, width: int = 900, height: int = 650):
    tab_ids = ["main", "format", "built", "next"]
    paths = []
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": width, "height": height})
        page.goto(f"file:///{html_path.as_posix()}")
        for idx, tab_id in enumerate(tab_ids, start=1):
            page.locator(f"[data-tab='{tab_id}']").click()
            page.wait_for_timeout(150)
            png_path = html_path.parent / f"{slug}_tab{idx}.png"
            page.screenshot(path=str(png_path), full_page=False)
            paths.append(str(png_path))
        browser.close()
    return paths


def screenshot_html(html_path: Path, png_path: Path, width: int = 900, height: int = 650):
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": width, "height": height})
        page.goto(f"file:///{html_path.as_posix()}")
        page.screenshot(path=str(png_path), full_page=False)
        browser.close()


def slugify(text: str) -> str:
    text = re.sub(r"[^\w\s-]", "", text).strip().lower()
    return re.sub(r"[-\s]+", "-", text) or "summary"


def main():
    parser = argparse.ArgumentParser(description="Render ClawVision summary visual.")
    parser.add_argument("--summary", "-s", help="Path to JSON summary or '-' for stdin")
    parser.add_argument("--output", "-o", required=True, help="Output directory")
    parser.add_argument("--slug", help="Output slug (default: derived from title)")
    parser.add_argument("--lang", default="en", help="Language code for UI labels")
    parser.add_argument("--png", action="store_true", help="Render PNG screenshots")
    parser.add_argument("--md", action="store_true", help="Render Markdown summary")
    parser.add_argument("--pptx", action="store_true", help="Render PowerPoint deck")
    args = parser.parse_args()

    if args.summary == "-" or args.summary is None:
        raw = sys.stdin.read()
    else:
        raw = Path(args.summary).read_text(encoding="utf-8")

    summary = json.loads(raw)

    out_dir = Path(args.output)
    out_dir.mkdir(parents=True, exist_ok=True)

    slug = args.slug or slugify(summary.get("title", "summary"))
    md_file = f"{slug}.md"
    pptx_file = f"{slug}.pptx"
    html_path = out_dir / f"{slug}.html"
    html_path.write_text(render_html(summary, slug, md_file, pptx_file, lang=args.lang), encoding="utf-8")

    result = {"html": str(html_path)}

    if args.md:
        md_path = out_dir / md_file
        md_path.write_text(render_md(summary, lang=args.lang), encoding="utf-8")
        result["md"] = str(md_path)

    if args.pptx:
        pptx_path = out_dir / pptx_file
        render_pptx(summary, pptx_path, lang=args.lang)
        result["pptx"] = str(pptx_path)

    if args.png:
        tab_paths = screenshot_tabs(html_path, slug)
        result["png_tabs"] = tab_paths
        png_path = out_dir / f"{slug}.png"
        screenshot_html(html_path, png_path)
        result["png"] = str(png_path)

    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
